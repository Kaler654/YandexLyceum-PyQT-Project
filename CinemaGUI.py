from CinemaClasses import Cinema, CinemaHall, CinemaNetwork
import sys
from PyQt5.QtWidgets import QApplication, QTableWidgetItem, QLabel, QPushButton, QMainWindow, \
    QInputDialog, QLineEdit
from PyQt5.QtCore import Qt
from PyQt5.QtGui import QFont
import docx
from pptx import Presentation
import csv
import sqlite3
from authorization import AuthorizationDesign
from user_start_menu import UserStartMenuDesign
from account_info import AccountInfoDesign
from developer_menu import DeveloperMenuDesign
from timetable import TimeTableDesign
from cinemas import CinemasDesign
from bookseats import BookSeatsDesign
from actionselectionwindow import ActionSelectionWindowDesign


class MovingWindow:
    def mousePressEvent(self, event):
        """Функция вызывается при нажатии кнопки мыши"""
        if event.button() == Qt.LeftButton:
            self.old_pos = event.pos()

    def mouseReleaseEvent(self, event):
        """Функция вызывается при отпускании кнопки мыши"""
        if event.button() == Qt.LeftButton:
            self.old_pos = None

    def mouseMoveEvent(self, event):
        """Функция вызывается всякий раз, когда мышь перемещается"""
        try:
            if not self.old_pos:
                return
            delta = event.pos() - self.old_pos
            self.move(self.pos() + delta)
        except AttributeError:
            ""


class Authorization(MovingWindow, QMainWindow, AuthorizationDesign):
    def __init__(self):
        super(Authorization, self).__init__()
        self.setupUi(self)
        self.initUI()

    def initUI(self):
        self.setWindowFlags(Qt.FramelessWindowHint)
        self.setWindowTitle('Авторизация')
        self.setFixedSize(600, 400)
        self.pushButton_2.clicked.connect(self.registration)
        self.pushButton_3.clicked.connect(self.close_window)
        self.pushButton_4.clicked.connect(self.sign_in)
        self.lineEdit_2.setEchoMode(QLineEdit.Password)
        self.label_4.hide()

    def check_login(self, login):
        """Проверка есть ли логин в базе данных"""
        con = sqlite3.connect("users.sqlite")
        cur = con.cursor()
        result = cur.execute(f"""SELECT login FROM users
                                WHERE login = '{login}'""").fetchall()
        if len(result) == 0:
            return True
        return False

    def sign_in(self):
        """При правильных данных открывает меню посетителя"""
        con = sqlite3.connect("users.sqlite")
        cur = con.cursor()
        try:
            result = cur.execute(f"""SELECT login, password, money FROM users 
                                    WHERE login = '{self.lineEdit.text()}' AND 
                                    password = '{self.lineEdit_2.text()}'""").fetchall()[0]
        except IndexError:
            result = tuple()
        if len(result) == 0:
            self.label_4.setText("Неккоректный логин и/или пароль")
            self.label_4.show()
        elif result[0] == "admin":
            self.user_start_menu = UserStartMenu(result)
            self.developer_menu = DeveloperMenu(self.user_start_menu)
            self.developer_menu.show()
            self.close()
        else:
            self.user_start_menu = UserStartMenu(result)
            self.user_start_menu.show()
            self.close()

        con.close()

    def registration(self):
        """Регистрация нового пользователя"""
        self.label_4.hide()
        login, ok_pressed = QInputDialog.getText(QMainWindow(), "Введите информацию", "Логин")
        if ok_pressed and self.check_login(login):
            password, ok_pressed = QInputDialog.getText(QMainWindow(), "Введите информацию",
                                                        "Пароль")
            if ok_pressed:
                con = sqlite3.connect("users.sqlite")
                cur = con.cursor()
                result = cur.execute(f"""INSERT INTO users(login,password, money)
                                        VALUES('{login}','{password}',0)""")
                con.commit()
                con.close()
        elif ok_pressed and not self.check_login(login):
            self.label_4.setText("Пользователь с таким логином уже существует")
            self.label_4.show()

    def close_window(self):
        """Закрывает окно"""
        self.close()


class UserStartMenu(MovingWindow, QMainWindow, UserStartMenuDesign):
    def __init__(self, user_info):
        super(UserStartMenu, self).__init__()
        self.setupUi(self)
        self.user_info = user_info
        self.initUI(user_info)

    def initUI(self, user_info):
        self.setWindowFlags(Qt.FramelessWindowHint)
        self.setWindowTitle('Меню')
        self.setFixedSize(620, 200)
        self.pushButton_4.clicked.connect(self.buttons_handling)
        self.pushButton_5.clicked.connect(self.buttons_handling)
        self.pushButton_6.clicked.connect(self.buttons_handling)

    def buttons_handling(self):
        """Обработка нажатий кнопок"""
        sender = self.sender()
        if sender == self.pushButton_4:
            self.cinemas = Cinemas(self.user_info, self)
            self.cinemas.show()
            self.hide()
        elif sender == self.pushButton_5:
            self.account_info = AccountInfo(self.user_info, self)
            self.account_info.show()
            self.hide()
        elif sender == self.pushButton_6:
            if self.user_info[0] != "admin":
                self.close()
            else:
                self.developer_menu = DeveloperMenu(self)
                self.developer_menu.show()
                self.close()


class AccountInfo(MovingWindow, QMainWindow, AccountInfoDesign):
    def __init__(self, user_info, previous_window):
        super(AccountInfo, self).__init__()
        self.setupUi(self)
        self.login, self.password, self.balance = user_info
        self.previous_window = previous_window
        self.initUI(user_info)
        self.pushButton_3.clicked.connect(self.add_money_to_balance)
        self.pushButton_6.clicked.connect(self.close_window)
        self.pushButton_7.clicked.connect(self.show_password)
        self.filling_information()
        self.password_visible = False

    def initUI(self, user_info):
        self.setWindowFlags(Qt.FramelessWindowHint)
        self.setWindowTitle('Информация об аккаунте')
        self.setFixedSize(400, 270)

    def add_money_to_balance(self):
        """Добавление денег на баланс"""
        money, ok_pressed = QInputDialog.getInt(
            QMainWindow(), "Введите информацию", "Сумма пополнения баланса",
            100, 100, 99999, 1)
        if ok_pressed:
            con = sqlite3.connect("users.sqlite")
            cur = con.cursor()
            result = cur.execute(f"""UPDATE users
                                    SET money = {self.get_balance() + money}
                                    WHERE login = '{self.login}'""")
            con.commit()
            con.close()
            self.label_2.setText(f"Баланс: {self.get_balance()}")

    def show_password(self):
        """Показывает или скрывает пароль"""
        if self.password_visible:
            self.label_5.setText(f"Пароль: {len(self.password) * '*'}")
            self.password_visible = False
        else:
            self.label_5.setText(f"Пароль: {self.password}")
            self.password_visible = True

    def filling_information(self):
        """Заполнение информации о пользователе"""
        if self.login == "admin":
            account_type = "Разработчик"
        else:
            account_type = "Посетитель"
        self.label_3.setText(f"Тип аккаунта: {account_type}")
        self.label_2.setText(f"Баланс: {self.get_balance()}")
        self.label_4.setText(f"Логин: {self.login}")
        self.label_5.setText(f"Пароль: {len(self.password) * '*'}")

    def close_window(self):
        """Закрытие окна"""
        self.previous_window.show()
        self.close()

    def get_balance(self):
        """Получить баланс пользователя"""
        con = sqlite3.connect("users.sqlite")
        cur = con.cursor()
        result = cur.execute(f"""SELECT login, password, money FROM users 
                                 WHERE login = '{self.login}' AND 
                                 password = '{self.password}'""").fetchall()[0]
        return result[2]


class DeveloperMenu(MovingWindow, QMainWindow, DeveloperMenuDesign):
    def __init__(self, user_start_menu):
        super(DeveloperMenu, self).__init__()
        self.setupUi(self)
        self.initUI(user_start_menu)
        self.pushButton.clicked.connect(self.add_cinema_button)
        self.pushButton_2.clicked.connect(self.add_cinema_hall_button)
        self.pushButton_3.clicked.connect(self.change_seat_configuration)
        self.pushButton_4.clicked.connect(self.change_film)
        self.pushButton_5.clicked.connect(self.delete_cinema_in_cinema_network)
        self.pushButton_6.clicked.connect(self.workload_cinemas)
        self.pushButton_7.clicked.connect(self.brochures)
        self.pushButton_8.clicked.connect(self.timetable)
        self.pushButton_9.clicked.connect(self.buttons_handling)
        self.pushButton_10.clicked.connect(self.buttons_handling)

    def initUI(self, user_start_menu):
        self.setWindowFlags(Qt.FramelessWindowHint)
        self.setFixedSize(800, 500)
        self.setWindowTitle('Меню разработчика')
        self.user_start_menu = user_start_menu

    def add_cinema_button(self):
        """Функция добавляет новый кинотеатр при нажатии кнопки"""
        cinema_name, ok_pressed = QInputDialog.getText(QMainWindow(), "Введите информацию",
                                                       "Введите название кинотеатра")
        if ok_pressed:
            cinema_network.add_cinema(cinema_name)
            self.label_2.setText("Вы успешно добавили кинотеатр")
        else:
            self.label_2.setText("Не удалось добавить кинотеатр")

    def add_cinema_hall_button(self):
        """Функция добавляет новый зал по нажатию кнопки"""
        cinema_name, ok_pressed = \
            QInputDialog.getItem(QMainWindow(),
                                 "Введите информацию",
                                 "Выберите кинотеатр в которой хотите добавить зал",
                                 (cinema_network.get_cinema_names()), 0, False)
        if ok_pressed:
            hall_name, ok_pressed = QInputDialog.getText(QMainWindow(), "Введите кинотеатр",
                                                         "Введите название для нового зала")
            if ok_pressed:
                movie, ok_pressed = \
                    QInputDialog.getText(QMainWindow(), "Введите информацию",
                                         "Введите название для фильма в этом зале")
                if ok_pressed:
                    movie_time, ok_pressed = \
                        QInputDialog.getText(QMainWindow(), "Введите информацию",
                                             "Введите время начала фильма(в формате ЧЧ:ММ)")
                    if ok_pressed:
                        movie_length, ok_pressed = QInputDialog.getInt(
                            QMainWindow(), "Введите информацию", "Длительность фильма в минутах",
                            60, 1, 300, 1)
                        if ok_pressed:
                            price, ok_pressed = QInputDialog.getInt(
                                QMainWindow(), "Введите информацию", "Стоимость билета на фильм",
                                250, 100, 1000, 1)
                            if ok_pressed:
                                row, ok_pressed = QInputDialog.getInt(
                                    QMainWindow(), "Введите информацию", "Укажите число рядов",
                                    7, 1, 7, 1)
                                if ok_pressed:
                                    col, ok_pressed = \
                                        row, ok_pressed = QInputDialog.getInt(
                                        QMainWindow(),
                                        "Введите информацию", "Укажите число мест в ряду",
                                        14, 1, 14, 1)
                                    if ok_pressed:
                                        check = cinema_network.get_cinema_with_want_name(
                                            cinema_name).add_cinema_hall(hall_name, movie,
                                                                         movie_time, movie_length,
                                                                         price, int(row), int(col))
                                        if not check:
                                            self.label_2.setText("В кинотеатр "
                                                                 "нельзя добавить "
                                                                 "новый зал, так как он переполнен")
                                        else:
                                            self.label_2.setText("Вы успешно добавили новый зал")

    def change_seat_configuration(self):
        """Изменение конфигурации кресел"""
        cinema_name, ok_pressed = \
            QInputDialog.getItem(QMainWindow(),
                                 "Введите информацию",
                                 "Выберите название кинотетра в зале "
                                 "которого хотите изменить конфигурацию кресел",
                                 (cinema_network.get_cinema_names()), 0, False)

        if ok_pressed:
            hall_name, ok_pressed = \
                QInputDialog.getItem(QMainWindow(),
                                     "Введите информацию",
                                     "Выберите название зала в котором хотите "
                                     "изменить конфигруацию кресел",
                                     (cinema_network.get_cinema_with_want_name(
                                         cinema_name).get_cinema_hall_names()), 0, False)
            row, ok_pressed = QInputDialog.getItem(QMainWindow(),
                                                   "Выберите кол-во рядов",
                                                   "Выберите кол-во рядов", (
                                                       "1", "2", "3", "4", "5",
                                                       "6",
                                                       "7"), 0, False)
            if ok_pressed:
                col, ok_pressed = \
                    QInputDialog.getItem(QMainWindow(),
                                         "Введите информацию",
                                         "Выберите кол-во мест в ряду",
                                         ("1", "2", "3", "4", "5", "6", "7",
                                          "8", "9", "10", "11", "12", "13",
                                          "14"), 0, False)
                if ok_pressed:
                    cinema_network.get_cinema_with_want_name(cinema_name).change_seats(hall_name,
                                                                                       int(row),
                                                                                       int(col))
                    self.label_2.setText("Вы успешно изменили конфигурацию кресел")

    def change_film(self):
        """Функция изменяет фильм в кинозале"""
        cinema_name, ok_pressed = \
            QInputDialog.getItem(QMainWindow(),
                                 "Введите информацию",
                                 "Выберите название кинотетра в зале "
                                 "которого хотите изменить фильм",
                                 (cinema_network.get_cinema_names()), 0, False)
        if ok_pressed:
            hall_name, ok_pressed = \
                QInputDialog.getItem(QMainWindow(),
                                     "Введите информацию",
                                     "Выберите название зала в котором хотите "
                                     "изменить фильм",
                                     (cinema_network.get_cinema_with_want_name(
                                         cinema_name).get_cinema_hall_names()), 0, False)
            if ok_pressed:
                movie, ok_pressed = \
                    QInputDialog.getText(QMainWindow(), "Введите информацию",
                                         "Введите новое название для фильма в этом зале")
                if ok_pressed:
                    movie_time, ok_pressed = \
                        QInputDialog.getText(QMainWindow(), "Введите информацию",
                                             "Введите время начала нового фильма(в формате ЧЧ:ММ)")
                    if ok_pressed:
                        movie_length, ok_pressed = \
                            QInputDialog.getText(QMainWindow(), "Введите информацию",
                                                 "Введите продолжительность нового фильма в минутах")
                        if ok_pressed:
                            cinema_network.get_cinema_with_want_name(
                                cinema_name).get_hall_with_want_name(hall_name).replace_film(
                                movie, movie_time, movie_length)

    def delete_cinema_in_cinema_network(self):
        """Удаление кинотеатра из сети кинотеатров"""
        cinema_name, ok_pressed = \
            QInputDialog.getItem(QMainWindow(),
                                 "Введите информацию",
                                 "Выберите кинотеатр который хотите удалить",
                                 (cinema_network.get_cinema_names()), 0, False)
        cinema_network.delete_cinema(cinema_name)

    def workload_cinemas(self):
        """Функция выгружает загруженность кинотеатров"""
        records = []
        for cinema in cinema_network.cinemas:
            for hall in cinema.cinema_halls:
                records.append([hall.movie, hall.movie_time, hall.get_occupied_seat()])
        doc = docx.Document()

        # добавляем таблицу 3x3
        table = doc.add_table(rows=len(records) - (len(records) - 1), cols=3)
        # применяем стиль для таблицы
        table.style = 'Table Grid'
        hdr_cells = table.rows[0].cells
        hdr_cells[0].text = "Фильм"
        hdr_cells[1].text = "Время начала"
        hdr_cells[2].text = "Занятых мест"
        # заполняем таблицу данными
        for movie, time, seats in records:
            row_Cells = table.add_row().cells
            row_Cells[0].text = str(movie)
            row_Cells[1].text = str(time)
            row_Cells[2].text = str(seats)
        doc.save('Загруженность.docx')
        self.label_2.setText("Выгрузка загруженности прошла успешно")

    def brochures(self):
        """Рекламные буклеты о фильмах в кинотеатре"""
        cinema_name, ok_pressed = \
            QInputDialog.getItem(QMainWindow(),
                                 "Введите информацию",
                                 "Выберите название кинотеатра фильмы которого нужно выгрузить",
                                 (cinema_network.get_cinema_names()), 0, False)
        if ok_pressed:
            cinema_hall_names = []
            movies = []
            movies_time = []
            for hall in cinema_network.get_cinema_with_want_name(cinema_name).cinema_halls:
                cinema_hall_names.append(hall.hall_name)
                movies.append(hall.movie)
                movies_time.append(hall.movie_time)

            prs = Presentation()

            layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(layout)
            title = slide.shapes.title
            title.text = f"Фильмы в кинотеатре {cinema_name}"
            for i in range(0, len(cinema_hall_names)):
                layout = prs.slide_layouts[0]
                slide = prs.slides.add_slide(layout)
                title = slide.shapes.title
                title.text = f'Фильм "{movies[i]}" в зале "{cinema_hall_names[i]}"'
                subtitle = slide.placeholders[1]
                if movies[i] == "Начало":
                    subtitle.text = "Кобб – талантливый вор, лучший из " \
                                    "лучших в опасном искусстве извлечения: " \
                                    "он крадет ценные секреты из глубин " \
                                    "подсознания во время сна, когда человеческий " \
                                    "разум наиболее уязвим"
                elif movies[i] == "Мстители: Финал":
                    subtitle.text = "Оставшиеся в живых члены команды " \
                                    "Мстителей и их союзники должны разработать " \
                                    "новый план, который поможет противостоять " \
                                    "разрушительным действиям могущественного титана Таноса."
                elif movies[i] == "Король и Лев":
                    subtitle.text = "История об отважном львенке по имени Симба. " \
                                    "Знакомые с детства герои взрослеют, влюбляются, " \
                                    "познают себя и окружающий мир, совершают " \
                                    "ошибки и делают правильный выбор."
                elif movies[i] == "Интерстеллар":
                    subtitle.text = "Когда засуха, пыльные бури и вымирание растений " \
                                    "приводят человечество к продовольственному кризису, " \
                                    "коллектив исследователей и учёных отправляется " \
                                    "сквозь червоточину в путешествие."
                elif movies[i] == "Побег из Шоушенка":
                    subtitle.text = "Бухгалтер Энди Дюфрейн обвинён в убийстве " \
                                    "собственной жены и её любовника. Оказавшись в " \
                                    "тюрьме под названием Шоушенк, он сталкивается с " \
                                    "жестокостью и беззаконием, царящими по обе стороны " \
                                    "решётки."
                elif movies[i] == "Форрест Гамп":
                    subtitle.text = "От лица главного героя Форреста Гампа, слабоумного " \
                                    "безобидного человека с благородным и открытым " \
                                    "сердцем, рассказывается история его необыкновенной " \
                                    "жизни."
                elif movies[i] == "Список Шиндлера":
                    subtitle.text = "Фильм рассказывает реальную историю загадочного " \
                                    "Оскара Шиндлера, члена нацистской партии, " \
                                    "преуспевающего фабриканта, спасшего во время " \
                                    "Второй мировой войны почти 1200 евреев."
                elif movies[i] == "Бойцовский клуб":
                    subtitle.text = "Сотрудник страховой компании страдает " \
                                    "хронической бессонницей и отчаянно " \
                                    "пытается вырваться из мучительно скучной жизни."
                elif movies[i] == "Властелин колец":
                    subtitle.text = "Снятая режиссёром Питером Джексоном " \
                                    "серия из трёх связанных единым сюжетом " \
                                    "кинофильмов, представляющая собой " \
                                    "экранизацию романа Дж. Р. Р. Толкина «Властелин колец»."
                else:
                    subtitle.text = "Скоро в наших кинотеатрах..."
                prs.save("Фильмы кинотеатра.pptx")
                self.label_2.setText("Успешно выгружены рекламные буклеты")

    def timetable(self):
        """Расписание"""
        self.timetable = TimeTable(self)
        self.timetable.show()
        self.hide()

    def buttons_handling(self):
        """Функция обрабатывает кнопки для дальнейшего вызова нужной функции"""
        sender = self.sender()
        self.label_2.setText("")

        if sender == self.pushButton_9:
            self.close()

        elif sender == self.pushButton_10:
            self.label_2.setText("")
            self.user_start_menu.show()
            self.close()


class TimeTable(MovingWindow, QMainWindow, TimeTableDesign):
    def __init__(self, previous_window):
        super(TimeTable, self).__init__()
        self.setupUi(self)
        self.previous_window = previous_window
        self.initUI(previous_window)

    def initUI(self, previous_window):
        self.setWindowFlags(Qt.FramelessWindowHint)
        self.setWindowTitle('Расписание')
        self.setFixedSize(680, 500)
        self.pushButton_9.clicked.connect(self.buttons_handling)
        self.create_csv()
        self.load_table("Расписание сеансов за прошедший месяц.csv")

    def buttons_handling(self):
        """Обработчик кнопок"""
        sender = self.sender()
        if sender == self.pushButton_9:
            self.previous_window.show()
            self.close()

    def create_csv(self):
        """Создание расписания"""
        with open('Расписание сеансов за прошедший месяц.csv', 'w', newline='') as csvfile:
            writer = csv.writer(csvfile, delimiter=';', quotechar='"', quoting=csv.QUOTE_MINIMAL)
            writer.writerow(["Кинотеатр", "Кинозал", "Фильм", "Время начала",
                             "Продолжительность фильма", "Кол-во мест"])
            for cinema in cinema_network.cinemas:
                for cinema_hall in cinema.cinema_halls:
                    cinema_name, cinema_hall_name, movie, movie_time, movie_length, seats = \
                        cinema.cinema_name, cinema_hall.hall_name, cinema_hall.movie, \
                        cinema_hall.movie_time, cinema_hall.movie_length, \
                        cinema_hall.row * cinema_hall.col
                    writer.writerow(
                        [cinema_name, cinema_hall_name, movie, movie_time, movie_length, seats])

    def load_table(self, table_name):
        """Вывод расписания для просмотра"""
        with open(table_name) as csvfile:
            reader = csv.reader(csvfile, delimiter=';', quotechar='"')
            title = next(reader)
            self.tableWidget.setColumnCount(len(title))
            self.tableWidget.setHorizontalHeaderLabels(title)
            self.tableWidget.setRowCount(0)
            for i, row in enumerate(reader):
                self.tableWidget.setRowCount(
                    self.tableWidget.rowCount() + 1)
                for j, elem in enumerate(row):
                    self.tableWidget.setItem(
                        i, j, QTableWidgetItem(elem))
        self.tableWidget.resizeColumnsToContents()


class Cinemas(MovingWindow, QMainWindow, CinemasDesign):
    def __init__(self, user_info, previous_window):
        super(Cinemas, self).__init__()
        self.setupUi(self)
        self.previous_window = previous_window
        self.user_info = user_info
        self.initUI(user_info)

    def initUI(self, user_info):
        self.setWindowFlags(Qt.FramelessWindowHint)
        self.setWindowTitle('Фильмы')
        self.setFixedSize(610, 555)
        self.buttons_cinemas, self.buttons_films = [], []
        self.pushButton_3.clicked.connect(self.show_movies)
        self.create_buttons_cinemas()

    def create_buttons_cinemas(self):
        """Создает кнопки для выбора кинотеатра"""
        left, top = 40, 90
        for i in range(0, len(cinema_network.cinemas)):
            button_cinema = QPushButton(cinema_network.cinemas[i].cinema_name, self)
            button_cinema.clicked.connect(self.show_movies)
            button_cinema.resize(140, 50)
            button_cinema.move(left, top)
            button_cinema.setStyleSheet("""background-color: #FFC200
            ;color: #19072c;font-weight: 700;font-size: 16px;""")
            top += 80
            self.buttons_cinemas.append(button_cinema)

    def create_buttons_films(self, cinema):
        """Создает кнопки выбора фильма"""
        main_label_film = QLabel("Выберите фильм", self)
        main_label_film.setStyleSheet("""color: #FFC200;font-weight: 700;font-size: 16px;""")
        main_label_film.move(380, 37)
        main_label_film.resize(250, 50)
        main_label_film.show()
        left, top = 330, 90
        for i in range(0, len(cinema.cinema_halls)):
            button_film = QPushButton("|" + cinema.cinema_halls[i].movie + "|", self)
            button_film_time = QPushButton(cinema.cinema_halls[i].movie_time + " / " +
                                           str(cinema.cinema_halls[i].price) + " руб", self)
            button_film.resize(250, 50)
            button_film_time.resize(250, 25)
            button_film.move(left, top)
            button_film_time.move(left, top + 50)
            button_film.clicked.connect(self.show_movies)
            button_film.setStyleSheet("""color: #FFC200;font-weight: 700;font-size: 16px;""")
            button_film_time.setStyleSheet(
                """color: #FFC200;font-weight: 700;font-size: 16px;background-color: #19072c;""")
            button_film.show()
            button_film_time.show()
            top += 120
            self.buttons_films.append(button_film)
            self.buttons_films.append(button_film_time)

    def show_movies(self):
        """Обработчик сигналов от кнопок и открытие следующего окна"""
        sender = self.sender()
        if sender in self.buttons_cinemas:
            self.cinema = [_ for _ in cinema_network.cinemas if _.cinema_name == sender.text()][0]
            for _ in self.buttons_films:
                _.hide()
            self.create_buttons_films(self.cinema)
        elif sender in self.buttons_films:
            cinema_hall = [_ for _ in self.cinema.cinema_halls if _.movie ==
                           sender.text().replace("|", "")][0]
            self.book_seats = BookSeats(self, cinema_hall, self.user_info)
            self.hide()
            self.book_seats.show()
        else:
            self.close()
            self.previous_window.show()


class BookSeats(MovingWindow, QMainWindow, BookSeatsDesign):
    def __init__(self, cinema_window, cinema_hall, user_info):
        super(BookSeats, self).__init__()
        self.setupUi(self)
        self.login, self.password, self.balance = user_info
        self.row, self.col, self.cinema_hall, self.cinema_window, self.count_of_reservation_seats = \
            cinema_hall.row, cinema_hall.col, cinema_hall, cinema_window, 0
        self.initUI(cinema_hall)
        if cinema_hall.row > 7:
            cinema_hall.row = 7
        if cinema_hall.col > 14:
            cinema_hall.col = 14
        self.create_seats(cinema_hall.row, cinema_hall.col, cinema_hall)
        self.create_numbering(cinema_hall.row)

    def initUI(self, cinema_hall):
        self.setWindowFlags(Qt.FramelessWindowHint)
        self.setWindowTitle('Бронирование фильма')
        self.setFixedSize(565, 400)
        self.label.setText(f'Бронирование мест на фильм "{self.cinema_hall.movie}"')
        self.pushButton_3.clicked.connect(self.book_seats)
        self.pushButton_3.show()
        self.pushButton_4.clicked.connect(self.book_seats)
        self.seats_buttons = []
        self.label2 = QLabel("", self)
        self.label2.resize(500, 50)
        self.label2.setStyleSheet("color: #FFC200;font-weight: 700;")
        self.label2.setFont(QFont("MS Shell Dlg 2", 10, QFont.Bold))
        self.label2.show()

    def book_seats(self):
        """Закрытие окна с бронированием места"""
        sender = self.sender()
        if sender == self.pushButton_3:
            self.action_selection_window = \
                ActionSelectionWindow(self.cinema_window, self.count_of_reservation_seats)
            self.action_selection_window.show()
            self.close()
        elif sender == self.pushButton_4:
            self.close()
            self.cinema_window.show()

    def create_numbering(self, row):
        """Функция слздает нумерацию рядов"""
        top = 100
        for _ in range(1, row + 1):
            numbering = QLabel(str(_), self)
            numbering.setStyleSheet("""color: #FFC200;""")
            numbering.move(15, top)
            top += 37
            numbering.resize(10, 10)
            numbering.setFont(QFont("Times", 9, QFont.Bold))
            numbering.show()

    def create_seats(self, count_rows, count_cols, cinema_hall):
        """Функция создаёт кнопки/сиденья"""
        seat_count, left, top = 0, 35, 88
        for row in range(count_rows):
            for col in range(count_cols):
                seat_button = QPushButton(str(col + 1), self)
                seat = cinema_hall.seats_in_the_hall[row][col]
                if seat == '[ ]':
                    seat_button.clicked.connect(self.reservation_seat)
                    seat_button.setStyleSheet("""color: white;background-color: green;""")
                else:
                    seat_button.setStyleSheet("""color: white;background-color: red;""")
                seat_button.resize(35, 35)
                seat_button.move(left, top)
                self.seats_buttons.append(seat_button)
                left += 35
            left = 35
            top += 37

    def reservation_seat(self):
        """Бронирование места в зале по нажатию кнопки"""
        sender = self.sender()
        index_seat_in_hall = self.seats_buttons.index(sender)
        row, col = index_seat_in_hall // self.col + 1, int(sender.text())
        if self.update_balance(self.cinema_hall.price):
            self.cinema_hall.reservation_seat_in_the_hall(row, col)
            sender.setStyleSheet("""color: white;background-color: red;""")
            sender.disconnect()
            self.label2.setText(f"Бронирование прошло успешно: ряд {row} место {col} "
                                f"Ваш баланс - {self.get_balance()} руб")
            self.count_of_reservation_seats += 1
            self.label2.move(40, 345)
        else:
            self.label2.setText("Пополните баланс")
            self.label2.move(210, 345)

    def update_balance(self, price):
        """Обновление баланса пользователя"""
        con = sqlite3.connect("users.sqlite")
        cur = con.cursor()
        self.balance = self.get_balance()

        if self.balance - price >= 0:
            self.balance -= price
            result = cur.execute(f"""UPDATE users
                                     SET money = {self.balance}
                                     WHERE login = '{self.login}'""")
            con.commit()
            con.close()
            return True
        return False

    def get_balance(self):
        """Получение баланса пользователя"""
        con = sqlite3.connect("users.sqlite")
        cur = con.cursor()
        result = cur.execute(f"""SELECT login, password, money FROM users 
                                 WHERE login = '{self.login}' AND 
                                 password = '{self.password}'""").fetchall()[0]
        return result[2]


class ActionSelectionWindow(MovingWindow, QMainWindow, ActionSelectionWindowDesign):
    def __init__(self, cinema_window, count_of_reservation_seats):
        super(ActionSelectionWindow, self).__init__()
        self.setupUi(self)
        self.initUI(cinema_window, count_of_reservation_seats)

    def initUI(self, cinema_window, count_of_reservation_seats):
        self.setWindowFlags(Qt.FramelessWindowHint)
        self.setWindowTitle('Бронирование фильма')
        self.setFixedSize(520, 180)
        self.label.setText(
            f"Забронировано мест: {count_of_reservation_seats}.Ожидайте своего сеанса.")
        self.cinema_window = cinema_window
        self.pushButton.clicked.connect(self.open_next_window)
        self.pushButton_2.clicked.connect(self.open_next_window)

    def open_next_window(self):
        """Открытие предыдущего окна или закрытие этого"""
        sender = self.sender()
        if sender == self.pushButton:
            self.close()
            self.cinema_window.show()
        else:
            self.close()


def except_hook(cls, exception, traceback):
    sys.__excepthook__(cls, exception, traceback)


if __name__ == '__main__':
    app = QApplication(sys.argv)
    app.setStyle('Fusion')
    ex = Authorization()
    ex.show()
    cinema_network = CinemaNetwork()
    sys.excepthook = except_hook
    sys.exit(app.exec_())
